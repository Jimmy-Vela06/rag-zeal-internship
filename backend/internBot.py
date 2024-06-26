import os
import openai
from pinecone import Pinecone
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv
import uvicorn
import logging
import concurrent.futures
import asyncio
from functools import lru_cache
from fastapi.responses import JSONResponse
import json
import tiktoken
from tqdm.auto import tqdm
import re
import pandas as pd
from pdfminer.high_level import extract_text
import docx2txt
from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET
import base64
import sklearn

# pip install scikit-learn

from sklearn.metrics.pairwise import cosine_similarity

# Load environment variables
load_dotenv()

# Constants
GPT_MODEL = "gpt-4"
EMBEDDING_MODEL = "text-embedding-ada-002"
PATH = "./data"
INDEX = "idx"

# Initialize OpenAI and Pinecone
openai.api_key = os.getenv("OPENAI_API_KEY")
pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
ENV = os.getenv("ENV")

# Set up logging
logging.basicConfig(level=logging.INFO)

# FastAPI setup
app = FastAPI()

# Allow CORS for local development
origins = [
    "http://localhost",
    "http://localhost:4200",
    "http://localhost:4200/rag",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Helper functions
def num_tokens(text: str, model: str = "gpt-4") -> int:
    encoding = tiktoken.encoding_for_model(model)
    return len(encoding.encode(text))


def read_txt(path):
    with open(path, encoding="utf8") as f:
        return f.readlines()


def read_pdf(path):
    return extract_text(path).split("\n")


def read_docx(path):
    return docx2txt.process(path).split("\n")


def read_pptx(path):
    split = []
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                split.append(shape.text)
    return split


def read_xlsx(path):
    with zipfile.ZipFile(path, "r") as zip_ref:
        zip_ref.extractall(PATH)
    xml_file = os.path.join(PATH, "xl/sharedStrings.xml")
    tree = ET.parse(xml_file)
    root = tree.getroot()
    result = []
    for elem in root.iter():
        if elem.text:
            result.append(elem.text.strip())
    return result


def read_png(path):
    with open(path, "rb") as image_file:
        image = base64.b64encode(image_file.read()).decode("utf-8")
        messages = [
            {"role": "system", "content": "system message here"},
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Convert the following image into a text table.",
                    },
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{image}"},
                    },
                ],
            },
        ]
        functions = [
            {
                "name": "write_queries",
                "description": "Get information from the database",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "details": {
                            "type": "string",
                            "description": "The information needed from the database",
                        },
                    },
                },
            }
        ]
        result = ""
        for chunk in openai.ChatCompletion.create(
            model=GPT_MODEL,
            messages=messages,
            functions=functions,
            temperature=0,
            frequency_penalty=0,
            presence_penalty=0,
            stream=True,
        ):
            if len(chunk.choices) > 0:
                response = chunk.choices[0]
                finish_reason = response.finish_reason
                content = response.delta.content
                function = response.delta.function_call
                if content is not None:
                    result += content
        return result.split("\n")


def collect_all_files(root_path):
    all_files = []
    for root, dirs, files in os.walk(root_path):
        for file in files:
            file_path = os.path.join(root, file)
            all_files.append(file_path)
    return all_files


def get_suffix(s, delimiter="."):
    parts = s.split(delimiter)
    return delimiter.join(parts[-1:])


@lru_cache(maxsize=128)
def cached_process_file(file):
    suffix = get_suffix(file)
    if suffix == "txt":
        return read_txt(file)
    elif suffix == "pdf":
        return read_pdf(file)
    elif suffix == "docx":
        return read_docx(file)
    elif suffix == "pptx":
        return read_pptx(file)
    elif suffix == "xlsx":
        return read_xlsx(file)
    elif suffix == "png":
        return read_png(file)
    return []


# Async function to process files
async def process_file(file):
    return cached_process_file(file)


# Collect and process files in parallel
async def process_all_files(root_path):
    files = collect_all_files(root_path)
    with concurrent.futures.ThreadPoolExecutor() as executor:
        loop = asyncio.get_event_loop()
        tasks = [loop.run_in_executor(executor, process_file, file) for file in files]
        results = await asyncio.gather(*tasks)
    return results


# FastAPI endpoint
class ChatRequest(BaseModel):
    message: str


@app.post("/chat")
async def chat(request: ChatRequest):
    chatbox = request.message

    if chatbox == "exit":
        return JSONResponse(content={"response": "Goodbye!"})

    # Embed the query
    embedding = (
        openai.Embedding.create(model=EMBEDDING_MODEL, input=chatbox).data[0].embedding
    )

    # Query Pinecone index
    if INDEX in [index.name for index in pc.list_indexes()]:
        idx = pc.Index(INDEX)
        result = idx.query(vector=[embedding], top_k=10, include_metadata=True)
        context = [x["metadata"]["text"].replace("\n", "") for x in result["matches"]]

        # Generate response using context
        messages = [
            {
                "role": "system",
                "content": "You are William Shakespeare. Keep sentences summarized to less than 20 words. If asked to explain, give full detail. Only respond using the context provided. Do not add any extra information. For any unrelated questions, tell the user to ask something about the context.",
            },
            {"role": "user", "content": chatbox},
            {"role": "system", "content": " ".join(context)},
        ]

        # Generate response
        response = openai.ChatCompletion.create(
            model="gpt-4", messages=messages, max_tokens=500, temperature=0.7
        )

        # Extract response text and strip unnecessary "/" "n" and spaces
        response_text = response.choices[0].message["content"].strip()
        response_text = response_text.replace("\\n", "\n")
        response_text = response_text.replace('\\"', '"')

        response_dict = {"response": response_text}
        return JSONResponse(content=response_dict)

    else:
        return JSONResponse(content={"response": "Index not found"})


class QueryRequest(BaseModel):
    query: str
    similar_vectors: str
    response_len: str
    temp: str
    perspective: str


class QueryResponse(BaseModel):
    results: list


async def get_embedding(query: str):
    response = await openai.Embedding.acreate(model=EMBEDDING_MODEL, input=query)
    return response["data"][0]["embedding"]


async def query_pinecone(embedding, top_k: int):
    return index.query(vector=embedding, top_k=top_k, include_metadata=True)


async def get_chat_completion(messages, max_tokens: int, temperature: float):
    return await openai.ChatCompletion.acreate(
        model="gpt-4", messages=messages, max_tokens=max_tokens, temperature=temperature
    )


@app.post("/query")
async def handle_query(request: QueryRequest):
    try:
        # Ensure Pinecone index is initialized
        if INDEX not in [index.name for index in pc.list_indexes()]:
            raise HTTPException(status_code=500, detail="Index not found")

        idx = pc.Index(INDEX)

        embedding = await get_embedding(request.query)
        pinecone_result = idx.query(
            vector=embedding, top_k=int(request.similar_vectors), include_metadata=True
        )

        context = [
            x["metadata"]["text"].replace("\n", "") for x in pinecone_result["matches"]
        ]
        context = context[:5]  # Limit to top 5 results

        messages_with_context = [
            {"role": "system", "content": request.perspective},
            {"role": "user", "content": request.query},
            {"role": "assistant", "content": "Based on the following context:"},
            {"role": "system", "content": " ".join(context)},
        ]

        messages_without_context = [
            {"role": "system", "content": request.perspective},
            {"role": "user", "content": request.query},
        ]

        response_with_context, response_without_context = await asyncio.gather(
            get_chat_completion(
                messages_with_context, int(request.response_len), float(request.temp)
            ),
            get_chat_completion(
                messages_without_context, int(request.response_len), float(request.temp)
            ),
        )

        return QueryResponse(
            results=[
                {
                    "context": context,
                    "response": response_with_context["choices"][0]["message"][
                        "content"
                    ].strip(),
                },
                {
                    "context": [],
                    "response": response_without_context["choices"][0]["message"][
                        "content"
                    ].strip(),
                },
            ]
        )
    except Exception as e:
        logging.error(f"Error processing query: {str(e)}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.get("/")
def read_root():
    return {"Hello": "World"}


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8001))
    uvicorn.run(app, host="0.0.0.0", port=port)
